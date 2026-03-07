import streamlit as st
import easyocr
import numpy as np
from PIL import Image, ImageDraw
import fitz
from docx import Document
from io import BytesIO
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Page Setup ---
st.set_page_config(page_title="AI Scanner Ultimate", layout="wide")
st.title("🚀 AI Scanner Ultimate (Word, Excel, PPT, JPG)")

# --- Sidebar ---
langs = st.sidebar.multiselect("Zaban (Languages):", ["en", "ur", "ar"], default=["en", "ur"])
zoom = st.sidebar.slider("Quality (2.0 is Best):", 1.0, 3.0, 2.0)

@st.cache_resource
def get_reader(l):
    return easyocr.Reader(l, gpu=False)

up_file = st.file_uploader("File select karein", type=["pdf", "png", "jpg", "jpeg"])

if up_file:
    reader = get_reader(langs)
    text_data = []

    if st.button("Scanning Shuru Karein"):
        with st.spinner("AI processing ho rahi hai..."):
            try:
                images_list = []
                if up_file.type == "application/pdf":
                    doc = fitz.open(stream=up_file.read(), filetype="pdf")
                    for page in doc:
                        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        images_list.append(img)
                        res = reader.readtext(np.array(img), detail=0)
                        text_data.extend(res)
                else:
                    img = Image.open(up_file)
                    images_list.append(img)
                    res = reader.readtext(np.array(img), detail=0)
                    text_data = res

                full_string = "\n".join(text_data)
                st.success("Scanning Mukammal!")

                # --- 1. WORD DOWNLOAD ---
                doc_file = Document()
                for line in text_data:
                    p = doc_file.add_paragraph(line)
                    if any(ord(c) > 1200 for c in line):
                        p.paragraph_format.alignment = 2 # Right
                    else: p.paragraph_format.alignment = 0
                
                word_buf = BytesIO()
                doc_file.save(word_buf)
                st.download_button("📥 Word Download", word_buf.getvalue(), "result.docx")

                # --- 2. EXCEL DOWNLOAD ---
                df = pd.DataFrame(text_data, columns=["Scanned Text"])
                excel_buf = BytesIO()
                with pd.ExcelWriter(excel_buf, engine='xlsxwriter') as writer:
                    df.to_excel(writer, index=False, sheet_name='Sheet1')
                st.download_button("📊 Excel Download", excel_buf.getvalue(), "result.xlsx")

                # --- 3. POWERPOINT DOWNLOAD ---
                ppt = Presentation()
                slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Scanned Data"
                content = slide.placeholders[1]
                content.text = full_string[:1000] # Pehla kuch data slide par
                
                ppt_buf = BytesIO()
                ppt.save(ppt_buf)
                st.download_button("📽️ PPT Download", ppt_buf.getvalue(), "result.pptx")

                # --- 4. JPG DOWNLOAD ---
                if images_list:
                    img_buf = BytesIO()
                    images_list[0].save(img_buf, format="JPEG")
                    st.download_button("🖼️ JPG Download (Page 1)", img_buf.getvalue(), "scanned.jpg")

                st.text_area("Live Preview", full_string, height=300)

            except Exception as e:
                st.error(f"Error: {e}")
